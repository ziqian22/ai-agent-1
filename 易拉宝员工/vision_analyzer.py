"""
Claude Vision分析器
使用Claude Opus 4.7的Vision能力分析产品图片和文档
"""

import anthropic
import base64
import json
from typing import Dict, List, Optional, Any, Tuple
from pathlib import Path
import io

# 文档处理库
try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from PIL import Image


class VisionAnalyzer:
    """Claude Vision分析器，用于提取产品信息"""

    def __init__(self, api_key: str, base_url: Optional[str] = None):
        """
        初始化Vision分析器

        Args:
            api_key: Claude API密钥
            base_url: API基础URL（用于中转服务）
        """
        if base_url:
            self.client = anthropic.Anthropic(
                api_key=api_key,
                base_url=base_url
            )
        else:
            self.client = anthropic.Anthropic(api_key=api_key)

        # 使用更通用的模型名称，兼容中转服务
        if base_url and ("zhouyang168" in base_url or "yunwu.ai" in base_url):
            # 中转服务使用的模型名称
            self.model = "claude-opus-4-7"
        else:
            # 官方API使用的模型名称
            self.model = "claude-opus-4-20250514"

    def analyze_image(self, image_path: str) -> Dict[str, Any]:
        """
        分析产品图片，提取结构化信息

        Args:
            image_path: 图片路径

        Returns:
            产品信息字典
        """
        try:
            # 读取图片
            with open(image_path, "rb") as f:
                image_data = base64.b64encode(f.read()).decode()

            # 判断图片类型
            file_ext = Path(image_path).suffix.lower()
            media_type_map = {
                '.png': 'image/png',
                '.jpg': 'image/jpeg',
                '.jpeg': 'image/jpeg',
                '.gif': 'image/gif',
                '.webp': 'image/webp'
            }
            media_type = media_type_map.get(file_ext, 'image/jpeg')

            # 构建Vision prompt
            prompt = self._build_analysis_prompt()

            # 调用Claude Opus 4.7 Vision API
            response = self.client.messages.create(
                model=self.model,
                max_tokens=4096,  # 增加token限制，确保能返回完整的产品特点列表
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": media_type,
                                "data": image_data
                            }
                        },
                        {
                            "type": "text",
                            "text": prompt
                        }
                    ]
                }]
            )

            # 解析响应
            result_text = response.content[0].text

            # 尝试提取JSON
            product_info = self._extract_json_from_response(result_text)

            return product_info

        except Exception as e:
            print(f"分析图片出错: {str(e)}")
            return {
                "error": str(e),
                "product_name": "",
                "brand": "",
                "slogan": "",
                "features": [],
                "product_type": "",
                "scenes": [],
                "colors": []
            }

    def _build_analysis_prompt(self) -> str:
        """构建分析提示词"""
        return """请仔细分析这张产品宣传图片，提取以下信息：

1. **产品名称**：完整的产品型号或名称
2. **品牌名称**：品牌或公司名称
3. **核心卖点**：主要的宣传语或口号（一句话）
4. **产品特点**：列出所有产品卖点和特性（每条独立，尽可能详细）
5. **产品类型**：判断产品类别（如：饮水机、化妆品、电子产品等）
6. **适用场景**：产品适合使用的场景（如：办公室、家庭、医院等）
7. **主要颜色**：产品的主要颜色

**重要：你必须只返回JSON格式，不要添加任何其他文字说明。**

返回格式：
{
  "product_name": "产品名称",
  "brand": "品牌名称",
  "slogan": "核心卖点",
  "features": ["特点1", "特点2", "特点3"],
  "product_type": "产品类型",
  "scenes": ["场景1", "场景2"],
  "colors": ["颜色1", "颜色2"]
}

请直接返回JSON，不要使用markdown代码块包裹。"""

    def _extract_json_from_response(self, text: str) -> Dict[str, Any]:
        """
        从响应文本中提取JSON

        Args:
            text: 响应文本

        Returns:
            解析后的字典
        """
        try:
            # 尝试直接解析
            return json.loads(text)
        except json.JSONDecodeError:
            # 尝试提取代码块中的JSON
            if "```json" in text:
                start = text.find("```json") + 7
                end = text.find("```", start)
                json_text = text[start:end].strip()
                return json.loads(json_text)
            elif "```" in text:
                start = text.find("```") + 3
                end = text.find("```", start)
                json_text = text[start:end].strip()
                return json.loads(json_text)
            else:
                # 如果无法解析，返回默认结构
                print(f"无法解析JSON，原始响应: {text}")
                return {
                    "product_name": "",
                    "brand": "",
                    "slogan": "",
                    "features": [],
                    "product_type": "",
                    "scenes": [],
                    "colors": [],
                    "raw_response": text
                }

    def ask_clarification(self, product_info: Dict[str, Any], question: str) -> str:
        """
        向Claude询问澄清问题

        Args:
            product_info: 已提取的产品信息
            question: 需要澄清的问题

        Returns:
            Claude的回答
        """
        try:
            prompt = f"""基于以下产品信息：

{json.dumps(product_info, ensure_ascii=False, indent=2)}

用户的问题是：{question}

请简洁明确地回答用户的问题。"""

            response = self.client.messages.create(
                model=self.model,
                max_tokens=500,
                messages=[{
                    "role": "user",
                    "content": prompt
                }]
            )

            return response.content[0].text

        except Exception as e:
            print(f"询问澄清问题出错: {str(e)}")
            return f"抱歉，我无法回答这个问题。错误：{str(e)}"

    def generate_questions(self, product_info: Dict[str, Any]) -> List[str]:
        """
        根据提取的产品信息生成主动提问

        Args:
            product_info: 已提取的产品信息

        Returns:
            问题列表
        """
        questions = []

        # 检查信息完整性
        if not product_info.get('product_name'):
            questions.append("我没有识别到产品名称，请问产品的完整名称是什么？")

        if not product_info.get('brand'):
            questions.append("我没有识别到品牌信息，请问品牌名称是什么？")

        if not product_info.get('features') or len(product_info.get('features', [])) < 3:
            questions.append("我只识别到少量产品特点，是否还有其他重要特点需要补充？")

        if not product_info.get('scenes'):
            questions.append("请问这个产品主要适用于哪些场景？（如：办公室、家庭、医院等）")

        # 如果信息完整，询问确认
        if not questions:
            questions.append("以上信息是否准确？如需修改请告诉我。")

        return questions[:3]  # 最多返回3个问题

    def test_connection(self) -> bool:
        """
        测试API连接

        Returns:
            连接是否正常
        """
        try:
            response = self.client.messages.create(
                model=self.model,
                max_tokens=10,
                messages=[{
                    "role": "user",
                    "content": "Hello"
                }]
            )
            return True
        except Exception as e:
            print(f"连接测试失败: {str(e)}")
            return False

    def analyze_document(self, file_path: str) -> Dict[str, Any]:
        """
        分析文档（PDF/Word/PPT），提取产品信息和图片

        Args:
            file_path: 文档路径

        Returns:
            {
                'product_info': {...},  # 产品信息
                'images': [...]         # 提取的图片路径列表
            }
        """
        file_ext = Path(file_path).suffix.lower()

        if file_ext == '.pdf':
            return self._analyze_pdf(file_path)
        elif file_ext in ['.docx', '.doc']:
            return self._analyze_word(file_path)
        elif file_ext in ['.pptx', '.ppt']:
            return self._analyze_ppt(file_path)
        else:
            return {
                'product_info': {
                    'error': f'不支持的文件格式: {file_ext}',
                    'product_name': '',
                    'brand': '',
                    'slogan': '',
                    'features': [],
                    'product_type': '',
                    'scenes': [],
                    'colors': []
                },
                'images': []
            }

    def _analyze_pdf(self, file_path: str) -> Dict[str, Any]:
        """
        分析 PDF 文档（使用 Claude 原生 PDF 支持）

        Args:
            file_path: PDF 文件路径

        Returns:
            分析结果字典
        """
        if not PDF_AVAILABLE:
            return {
                'product_info': {
                    'error': 'PyPDF2 未安装，无法处理 PDF 文件',
                    'product_name': '',
                    'brand': '',
                    'slogan': '',
                    'features': [],
                    'product_type': '',
                    'scenes': [],
                    'colors': []
                },
                'images': []
            }

        try:
            # 读取 PDF 文件
            with open(file_path, 'rb') as f:
                pdf_data = base64.b64encode(f.read()).decode()

            # 构建分析提示词
            prompt = self._build_document_analysis_prompt()

            # 调用 Claude API（使用 PDF 支持）
            response = self.client.messages.create(
                model=self.model,
                max_tokens=3000,
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "document",
                            "source": {
                                "type": "base64",
                                "media_type": "application/pdf",
                                "data": pdf_data
                            }
                        },
                        {
                            "type": "text",
                            "text": prompt
                        }
                    ]
                }]
            )

            # 解析响应
            result_text = response.content[0].text
            product_info = self._extract_json_from_response(result_text)

            # 尝试从 PDF 提取图片
            extracted_images = self._extract_images_from_pdf(file_path)

            return {
                'product_info': product_info,
                'images': extracted_images
            }

        except Exception as e:
            print(f"分析 PDF 出错: {str(e)}")
            return {
                'product_info': {
                    'error': str(e),
                    'product_name': '',
                    'brand': '',
                    'slogan': '',
                    'features': [],
                    'product_type': '',
                    'scenes': [],
                    'colors': []
                },
                'images': []
            }

    def _analyze_word(self, file_path: str) -> Dict[str, Any]:
        """
        分析 Word 文档

        Args:
            file_path: Word 文件路径

        Returns:
            分析结果字典
        """
        if not DOCX_AVAILABLE:
            return {
                'product_info': {
                    'error': 'python-docx 未安装，无法处理 Word 文件',
                    'product_name': '',
                    'brand': '',
                    'slogan': '',
                    'features': [],
                    'product_type': '',
                    'scenes': [],
                    'colors': []
                },
                'images': []
            }

        try:
            # 读取 Word 文档
            doc = Document(file_path)

            # 提取文本
            text_content = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())

            # 提取表格中的文本
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_content.append(cell.text.strip())

            full_text = '\n'.join(text_content)

            # 使用 Claude 分析文本
            prompt = f"""{self._build_document_analysis_prompt()}

文档内容：
{full_text[:5000]}  # 限制长度避免超过 token 限制
"""

            response = self.client.messages.create(
                model=self.model,
                max_tokens=3000,
                messages=[{
                    "role": "user",
                    "content": prompt
                }]
            )

            result_text = response.content[0].text
            product_info = self._extract_json_from_response(result_text)

            # 提取图片
            extracted_images = self._extract_images_from_word(file_path)

            return {
                'product_info': product_info,
                'images': extracted_images
            }

        except Exception as e:
            print(f"分析 Word 文档出错: {str(e)}")
            return {
                'product_info': {
                    'error': str(e),
                    'product_name': '',
                    'brand': '',
                    'slogan': '',
                    'features': [],
                    'product_type': '',
                    'scenes': [],
                    'colors': []
                },
                'images': []
            }

    def _analyze_ppt(self, file_path: str) -> Dict[str, Any]:
        """
        分析 PowerPoint 文档

        Args:
            file_path: PPT 文件路径

        Returns:
            分析结果字典
        """
        if not PPTX_AVAILABLE:
            return {
                'product_info': {
                    'error': 'python-pptx 未安装，无法处理 PPT 文件',
                    'product_name': '',
                    'brand': '',
                    'slogan': '',
                    'features': [],
                    'product_type': '',
                    'scenes': [],
                    'colors': []
                },
                'images': []
            }

        try:
            # 读取 PPT 文档
            prs = Presentation(file_path)

            # 提取文本
            text_content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())

            full_text = '\n'.join(text_content)

            # 使用 Claude 分析文本
            prompt = f"""{self._build_document_analysis_prompt()}

文档内容：
{full_text[:5000]}  # 限制长度避免超过 token 限制
"""

            response = self.client.messages.create(
                model=self.model,
                max_tokens=3000,
                messages=[{
                    "role": "user",
                    "content": prompt
                }]
            )

            result_text = response.content[0].text
            product_info = self._extract_json_from_response(result_text)

            # 提取图片
            extracted_images = self._extract_images_from_ppt(file_path)

            return {
                'product_info': product_info,
                'images': extracted_images
            }

        except Exception as e:
            print(f"分析 PPT 文档出错: {str(e)}")
            return {
                'product_info': {
                    'error': str(e),
                    'product_name': '',
                    'brand': '',
                    'slogan': '',
                    'features': [],
                    'product_type': '',
                    'scenes': [],
                    'colors': []
                },
                'images': []
            }

    def _build_document_analysis_prompt(self) -> str:
        """构建文档分析提示词"""
        return """请仔细分析这份产品文档，提取以下信息：

1. **产品名称**：完整的产品型号或名称
2. **品牌名称**：品牌或公司名称
3. **核心卖点**：主要的宣传语或口号（一句话）
4. **产品特点**：列出所有产品卖点和特性（每条独立，尽可能详细）
5. **产品类型**：判断产品类别（如：饮水机、化妆品、电子产品等）
6. **适用场景**：产品适合使用的场景（如：办公室、家庭、医院等）
7. **主要颜色**：产品的主要颜色（如果文档中有描述）

**重要：你必须只返回JSON格式，不要添加任何其他文字说明。**

返回格式：
{
  "product_name": "产品名称",
  "brand": "品牌名称",
  "slogan": "核心卖点",
  "features": ["特点1", "特点2", "特点3"],
  "product_type": "产品类型",
  "scenes": ["场景1", "场景2"],
  "colors": ["颜色1", "颜色2"]
}

请直接返回JSON，不要使用markdown代码块包裹。"""

    def _extract_images_from_pdf(self, file_path: str) -> List[str]:
        """从 PDF 提取图片"""
        if not PDF_AVAILABLE:
            return []

        extracted_images = []
        temp_dir = Path("temp_uploads") / "extracted_images"
        temp_dir.mkdir(parents=True, exist_ok=True)

        try:
            pdf_reader = PyPDF2.PdfReader(file_path)

            for page_num, page in enumerate(pdf_reader.pages):
                if '/XObject' in page['/Resources']:
                    xObject = page['/Resources']['/XObject'].get_object()

                    for obj_name in xObject:
                        obj = xObject[obj_name]

                        if obj['/Subtype'] == '/Image':
                            try:
                                # 提取图片数据
                                size = (obj['/Width'], obj['/Height'])
                                data = obj.get_data()

                                # 保存图片
                                image_path = temp_dir / f"pdf_page{page_num}_{obj_name[1:]}.png"

                                if obj['/ColorSpace'] == '/DeviceRGB':
                                    img = Image.frombytes('RGB', size, data)
                                elif obj['/ColorSpace'] == '/DeviceGray':
                                    img = Image.frombytes('L', size, data)
                                else:
                                    continue

                                img.save(image_path)
                                extracted_images.append(str(image_path))

                            except Exception as e:
                                print(f"提取 PDF 图片出错: {str(e)}")
                                continue

        except Exception as e:
            print(f"读取 PDF 出错: {str(e)}")

        return extracted_images

    def _extract_images_from_word(self, file_path: str) -> List[str]:
        """从 Word 文档提取图片"""
        if not DOCX_AVAILABLE:
            return []

        extracted_images = []
        temp_dir = Path("temp_uploads") / "extracted_images"
        temp_dir.mkdir(parents=True, exist_ok=True)

        try:
            doc = Document(file_path)

            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_data = rel.target_part.blob
                        image_ext = rel.target_ref.split('.')[-1]

                        image_path = temp_dir / f"word_{rel.rId}.{image_ext}"

                        with open(image_path, 'wb') as f:
                            f.write(image_data)

                        extracted_images.append(str(image_path))

                    except Exception as e:
                        print(f"提取 Word 图片出错: {str(e)}")
                        continue

        except Exception as e:
            print(f"读取 Word 文档出错: {str(e)}")

        return extracted_images

    def _extract_images_from_ppt(self, file_path: str) -> List[str]:
        """从 PowerPoint 提取图片"""
        if not PPTX_AVAILABLE:
            return []

        extracted_images = []
        temp_dir = Path("temp_uploads") / "extracted_images"
        temp_dir.mkdir(parents=True, exist_ok=True)

        try:
            prs = Presentation(file_path)

            image_count = 0
            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            image_ext = image.ext

                            image_path = temp_dir / f"ppt_slide{slide_num}_{image_count}.{image_ext}"

                            with open(image_path, 'wb') as f:
                                f.write(image_bytes)

                            extracted_images.append(str(image_path))
                            image_count += 1

                        except Exception as e:
                            print(f"提取 PPT 图片出错: {str(e)}")
                            continue

        except Exception as e:
            print(f"读取 PPT 文档出错: {str(e)}")

        return extracted_images



# 测试代码
if __name__ == "__main__":
    import os
    from dotenv import load_dotenv

    load_dotenv()

    # 从环境变量获取API密钥
    api_key = os.getenv("CLAUDE_API_KEY")
    base_url = os.getenv("CLAUDE_BASE_URL")  # 可选

    if not api_key:
        print("错误：未找到CLAUDE_API_KEY环境变量")
        exit(1)

    # 创建分析器
    analyzer = VisionAnalyzer(api_key, base_url)

    # 测试连接
    print("测试API连接...")
    if analyzer.test_connection():
        print("✅ API连接正常")
    else:
        print("❌ API连接失败")
        exit(1)

    # 测试图片分析
    test_image = "74f8c6f71010b7553325446bec13f5a5.jpg"
    if Path(test_image).exists():
        print(f"\n分析图片: {test_image}")
        result = analyzer.analyze_image(test_image)
        print("\n提取的产品信息:")
        print(json.dumps(result, ensure_ascii=False, indent=2))

        # 生成主动提问
        questions = analyzer.generate_questions(result)
        print("\n主动提问:")
        for i, q in enumerate(questions, 1):
            print(f"{i}. {q}")
    else:
        print(f"测试图片不存在: {test_image}")
